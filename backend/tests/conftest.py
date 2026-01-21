"""Pytest configuration and fixtures."""

import pytest
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt


@pytest.fixture
def create_test_pptx():
    """Factory fixture to create test PPTX files."""

    def _create_pptx(slides_content: list[dict]) -> bytes:
        """
        Create a test PPTX with specified content.

        Args:
            slides_content: List of dicts with keys:
                - texts: list of (text, x, y, w, h) tuples
                  where x, y, w, h are in inches

        Returns:
            PPTX file as bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        for slide_data in slides_content:
            slide = prs.slides.add_slide(blank_layout)

            for text_item in slide_data.get("texts", []):
                text, x, y, w, h = text_item
                shape = slide.shapes.add_textbox(
                    Inches(x), Inches(y), Inches(w), Inches(h)
                )
                tf = shape.text_frame
                tf.text = text

        # Save to bytes
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    return _create_pptx


@pytest.fixture
def sample_ethics_pptx(create_test_pptx):
    """Create a sample PPTX with ethics content."""
    slides = [
        {
            "texts": [
                ("Film Analizi: Etik Değerlendirme", 1, 1, 8, 1),
                ("1/5", 9.2, 0.2, 0.5, 0.3),  # Page number top-right
            ]
        },
        {
            "texts": [
                ("Film Adı: Schindler's List", 1, 1, 8, 0.5),
                ("Tür: Drama, Tarih", 1, 1.6, 8, 0.5),
                ("Yönetmen: Steven Spielberg", 1, 2.2, 8, 0.5),
                ("Senarist: Steven Zaillian", 1, 2.8, 8, 0.5),
                ("Oyuncular: Liam Neeson, Ben Kingsley", 1, 3.4, 8, 0.5),
                ("2/5", 9.2, 0.2, 0.5, 0.3),
            ]
        },
        {
            "texts": [
                ("Etik İlkeleri", 1, 0.8, 8, 0.5),
                ("Adalet: Filmde adalet kavramı işlenmektedir.", 1, 1.5, 8, 1),
                ("Özerklik: Bireylerin kendi kararlarını verme hakkı.", 1, 2.5, 8, 1),
                ("Yararseverlik: İyilik yapma ve fayda sağlama.", 1, 3.5, 8, 1),
                ("Zarar vermeme: Kötülük yapmamak.", 1, 4.5, 8, 1),
                ("Dürüstlük: Doğruluk ve şeffaflık.", 1, 5.5, 8, 1),
                ("3/5", 9.2, 0.2, 0.5, 0.3),
            ]
        },
        {
            "texts": [
                ("Sahne Analizi", 1, 0.8, 8, 0.5),
                ("Filmde Schindler'in fabrikasında çalışan Yahudileri kurtarma sahnesi.", 1, 1.5, 8, 2),
                ("Bu sahne adalet ve yararseverlik ilkelerini yansıtmaktadır.", 1, 4, 8, 1),
                ("4/5", 9.2, 0.2, 0.5, 0.3),
            ]
        },
        {
            "texts": [
                ("Düşünceler ve Değerlendirme", 1, 0.8, 8, 0.5),
                ("Film etik açıdan çok öğretici bir yapımdır.", 1, 1.5, 8, 2),
                ("5/5", 9.2, 0.2, 0.5, 0.3),
            ]
        },
    ]
    return create_test_pptx(slides)


@pytest.fixture
def minimal_pptx(create_test_pptx):
    """Create a minimal PPTX with just one slide."""
    slides = [
        {
            "texts": [
                ("Test Slide", 1, 1, 8, 1),
            ]
        }
    ]
    return create_test_pptx(slides)
