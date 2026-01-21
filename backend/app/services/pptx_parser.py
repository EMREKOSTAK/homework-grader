"""PPTX Parser Service - Extracts structured content from PowerPoint files."""

import re
import uuid
from io import BytesIO
from typing import Optional, List
from pptx import Presentation
from pptx.util import Emu
from pptx.shapes.base import BaseShape

from app.models.schemas import (
    TextElement,
    SlideData,
    DeckMeta,
    ParsedPPTX,
    BoundingBox,
    TextStyle,
)


class PPTXParser:
    """Parser for extracting structured content from PPTX files."""

    def __init__(self):
        self._whitespace_pattern = re.compile(r'\s+')

    def parse(self, file_content: bytes) -> ParsedPPTX:
        """
        Parse a PPTX file and extract structured content.

        Args:
            file_content: Raw bytes of the PPTX file

        Returns:
            ParsedPPTX object with all extracted content
        """
        prs = Presentation(BytesIO(file_content))

        # Extract deck metadata
        meta = DeckMeta(
            slide_width=float(prs.slide_width),
            slide_height=float(prs.slide_height),
            units="EMU",
            total_slides=len(prs.slides)
        )

        # Extract content from each slide
        slides = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_data = self._extract_slide_content(slide, slide_idx)
            slides.append(slide_data)

        return ParsedPPTX(meta=meta, slides=slides)

    def _extract_slide_content(self, slide, slide_no: int) -> SlideData:
        """Extract content from a single slide."""
        elements = []

        for shape in slide.shapes:
            text_elements = self._extract_text_from_shape(shape)
            elements.extend(text_elements)

        return SlideData(slide_no=slide_no, elements=elements)

    def _extract_text_from_shape(self, shape: BaseShape) -> List[TextElement]:
        """Extract text elements from a shape."""
        elements = []

        # Skip shapes without text frame
        if not shape.has_text_frame:
            return elements

        text_frame = shape.text_frame

        # Combine all paragraphs into full text for this shape
        full_text_parts = []
        font_size = None
        is_bold = None
        font_name = None

        for paragraph in text_frame.paragraphs:
            para_text_parts = []
            for run in paragraph.runs:
                if run.text:
                    para_text_parts.append(run.text)
                    # Capture style from first non-empty run
                    if font_size is None and run.font.size:
                        font_size = run.font.size.pt if run.font.size else None
                    if is_bold is None and run.font.bold is not None:
                        is_bold = run.font.bold
                    if font_name is None and run.font.name:
                        font_name = run.font.name

            if para_text_parts:
                full_text_parts.append(''.join(para_text_parts))

        if not full_text_parts:
            return elements

        raw_text = '\n'.join(full_text_parts)
        normalized_text = self._normalize_text(raw_text)

        if not normalized_text:
            return elements

        # Get bounding box
        bbox = BoundingBox(
            x=float(shape.left) if shape.left else 0.0,
            y=float(shape.top) if shape.top else 0.0,
            w=float(shape.width) if shape.width else 0.0,
            h=float(shape.height) if shape.height else 0.0,
        )

        # Create text style
        style = TextStyle(
            font_size=font_size,
            bold=is_bold,
            font_name=font_name,
        )

        element = TextElement(
            id=str(uuid.uuid4())[:8],
            type="text",
            text=normalized_text,
            raw_text=raw_text.strip(),
            bbox=bbox,
            style=style,
        )
        elements.append(element)

        return elements

    def _normalize_text(self, text: str) -> str:
        """
        Normalize text by trimming and collapsing whitespace.
        Preserves line breaks as single spaces.
        """
        if not text:
            return ""

        # Replace all whitespace sequences (including newlines) with single space
        normalized = self._whitespace_pattern.sub(' ', text)
        # Trim leading/trailing whitespace
        return normalized.strip()

    def get_text_in_region(
        self,
        slide_data: SlideData,
        deck_meta: DeckMeta,
        x_min_ratio: float,
        x_max_ratio: float,
        y_min_ratio: float,
        y_max_ratio: float,
    ) -> List[TextElement]:
        """
        Get text elements within a specified region of the slide.

        Args:
            slide_data: The slide to search
            deck_meta: Deck metadata for dimensions
            x_min_ratio, x_max_ratio: Horizontal bounds as ratio of slide width (0-1)
            y_min_ratio, y_max_ratio: Vertical bounds as ratio of slide height (0-1)

        Returns:
            List of text elements in the region
        """
        width = deck_meta.slide_width
        height = deck_meta.slide_height

        x_min = width * x_min_ratio
        x_max = width * x_max_ratio
        y_min = height * y_min_ratio
        y_max = height * y_max_ratio

        matching_elements = []

        for element in slide_data.elements:
            # Check if element's center is in the region
            elem_center_x = element.bbox.x + element.bbox.w / 2
            elem_center_y = element.bbox.y + element.bbox.h / 2

            if x_min <= elem_center_x <= x_max and y_min <= elem_center_y <= y_max:
                matching_elements.append(element)

        return matching_elements

    def search_text(
        self,
        parsed_pptx: ParsedPPTX,
        pattern: str,
        case_sensitive: bool = False,
    ) -> List[tuple[int, TextElement]]:
        """
        Search for text matching a pattern across all slides.

        Args:
            parsed_pptx: Parsed presentation
            pattern: Regex pattern to match
            case_sensitive: Whether to use case-sensitive matching

        Returns:
            List of (slide_no, element) tuples for matches
        """
        flags = 0 if case_sensitive else re.IGNORECASE
        regex = re.compile(pattern, flags)

        matches = []
        for slide in parsed_pptx.slides:
            for element in slide.elements:
                if regex.search(element.text):
                    matches.append((slide.slide_no, element))

        return matches
